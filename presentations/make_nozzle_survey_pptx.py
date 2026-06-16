# -*- coding: utf-8 -*-
"""超音速・極超音速ノズル設計技術 文献調査スライド生成スクリプト。

出典: .github/plans/nozzle-optimization-tool-survey.md (2026-06 調査)
1 論文 = 1 スライド。各スライドは「目的 / 方法 / 結果 / 得られた知見」で統一。
papers/ に PDF がある論文には代表図 (figs_final_*.png) を貼り込む。

実行:
    /home/sano/work/forge/.venv-pptx/bin/python presentations/make_nozzle_survey_pptx.py
出力:
    presentations/nozzle_design_survey.pptx
"""

import os
import struct
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
JP_FONT = "Meiryo UI"

# ---------------------------------------------------------------- 色・凡例
STATUS = {
    "◎": ("検証済 ◎", RGBColor(0x2E, 0x7D, 0x32)),
    "○": ("古典・標準 ○", RGBColor(0x15, 0x65, 0xC0)),
    "△": ("要追検証 △", RGBColor(0xEF, 0x6C, 0x00)),
}
GRAY = RGBColor(0x59, 0x59, 0x59)
DARK = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 4 区分のラベル色 (目的/方法/結果/知見)
FIELD_DEF = [
    ("目的", RGBColor(0x15, 0x65, 0xC0)),
    ("方法", RGBColor(0x00, 0x69, 0x5C)),
    ("結果", RGBColor(0xE6, 0x51, 0x00)),
    ("知見", RGBColor(0x6A, 0x1B, 0x9A)),
]
FIELD_KEYS = ["obj", "method", "result", "insight"]
FIELD_FULL = ["目的・背景", "方法", "結果", "得られた知見・本ツールへの含意"]

SECTIONS = {
    "S1": ("第1章  古典 — MOC 設計法の確立 (1947–1978)", RGBColor(0x1F, 0x4E, 0x79)),
    "S2": ("第2章  スロート曲率・収縮部の物理 (設計変数化の根拠)", RGBColor(0x00, 0x69, 0x5C)),
    "S3": ("第3章  NS-in-the-loop — 補正から内在化へ (1992–2025)", RGBColor(0xB7, 0x1C, 0x1C)),
    "S4": ("第4章  随伴法 — NS ループ内最適化の核心", RGBColor(0x4A, 0x14, 0x8C)),
    "S5": ("第5章  形状パラメータ化", RGBColor(0xE6, 0x51, 0x00)),
    "S6": ("第6章  サロゲート・多目的最適化 (パレートの担い手)", RGBColor(0x1B, 0x5E, 0x20)),
    "S7": ("第7章  参照アーキテクチャ (OSS 最適化基盤)", RGBColor(0x37, 0x47, 0x4F)),
    "S8": ("第8章  非平衡凝縮 — 極超音速風洞固有", RGBColor(0x02, 0x77, 0xBD)),
    "S9": ("第9章  DACS・SERN 固有設計", RGBColor(0x88, 0x0E, 0x4F)),
    "S10": ("第10章  ④デュアルベル — 高度補償ノズル (5機種版・追補)", RGBColor(0x00, 0x69, 0x79)),
}

# ---------------------------------------------------------------- 論文データ
# 各要素: sec, short, year, status, cite, obj, method, result, insight, (img)
P = []
def add(**kw): P.append(kw)

# ===== 第1章 古典 =====
add(sec="S1", short="遷音速スロート解の原点", year="1947", status="○",
    cite='R. Sauer, "General Characteristics of the Flow Through Nozzles at Near Critical Speeds," NACA TM-1147, 1947.',
    obj="スロート近傍の遷音速流れを解析的に記述し、MOC を始動する初期値線 (starting line) を与える。",
    method="遷音速微小擾乱方程式を級数展開で解く (一次近似)。",
    result="ソニックラインがスロートで平面でなく湾曲することを定式化。",
    insight="MOC starting line の理論的原点。後続の Hall (1962)・Kliegel–Levine (1969) 高次化の出発点。",
    img="figs_final_sauer1947.png", img_cap="出典: Sauer (1947) NACA TM-1147 Fig.1 — スロートで平面でなく湾曲する臨界 (ソニック) 線の定義図")
add(sec="S1", short="軸対称ラバルノズルの解析的設計", year="1949", status="○",
    cite='K. Foelsch, "The Analytical Design of an Axially Symmetric Laval Nozzle for a Parallel and Uniform Jet," J. Aeronautical Sciences, 16(3):161–166, 1949.',
    obj="平行・一様な超音速噴流を生む軸対称ラバルノズルを解析的に設計する。",
    method="放射状ソース流れと出口一様流を解析接続し、閉形式のコンタを導出。",
    result="MOC を陽に解かずに一様噴流コンタを得る公式を確立。",
    insight="「目標流れ → コンタ」設計思想の源流。現代でも初期形状生成 (ウォームスタート) として有効。")
add(sec="S1", short="推力最適ベルノズル (TOC) — 変分法 + MOC", year="1958", status="◎",
    cite='G. V. R. Rao, "Exhaust Nozzle Contour for Optimum Thrust," Jet Propulsion, 28(6):377–382, 1958. DOI 10.2514/8.7324',
    obj="与えられた長さ・膨張比のもとで推力が最大となるノズルコンタを求める。",
    method="変分法で control surface 上の必要流れ条件を導き、それを実現するコンタを MOC で構築。",
    result="15° 円錐ノズル比で約 60% 短縮する TOC/TOP ベルの基礎を確立。",
    insight="② DACS の軸長短縮の古典基盤。「下流円弧 Rd + MOC 旋回コンタ」の区分構成の標準。")
add(sec="S1", short="遷音速ノズル流れの級数解", year="1962", status="◎",
    cite='I. M. Hall, "Transonic Flow in Two-Dimensional and Axially-Symmetric Nozzles," Q. J. Mechanics and Applied Mathematics, 15(4):487–508, 1962.',
    obj="スロート遷音速解を Sauer より高次化し、starting line の精度を上げる。",
    method="無次元スロート壁曲率 1/R の冪級数で速度場の最初の 3 項を 2D・軸対称の両方で導出。",
    result="starting line がスロート壁曲率 R 1 個で一意にパラメータ化されることを確立。",
    insight="「スロート壁曲率 = 設計入力」の理論根拠。限界: 小曲率 R<1 で破綻 (流量係数が負)。")
add(sec="S1", short="小スロート曲率への遷音速解の拡張", year="1969", status="◎",
    cite='J. R. Kliegel & J. N. Levine, "Transonic Flow in Small Throat Radius of Curvature Nozzles," AIAA Journal, 7(7):1375–1378, 1969. DOI 10.2514/3.5355',
    obj="Hall 級数が破綻する小スロート曲率域でも有効な遷音速解を得る。",
    method="級数を 1/(R+1) の冪級数に再定式化する。",
    result="全 R で収束し、大 R では Hall に一致、R<1 のロケット級でも有効。",
    insight="本ツールの遷音速 starting line 生成器の推奨実装。小曲率スロートの探索を可能化。")
add(sec="S1", short="極超音速風洞ノズル設計の世界標準 (CONTUR)", year="1970/78", status="◎",
    cite='J. C. Sivells, "Aerodynamic Design of Axisymmetric Hypersonic Wind-Tunnel Nozzles," J. Spacecraft & Rockets, 7(11):1292–1299, 1970 / AEDC-TR-78-63 (プログラム CONTUR), 1978.',
    obj="出口で一様・平行流となる極超音速風洞ノズルを系統的に設計する。",
    method="目標中心線マッハ分布 → 軸対称 MOC → 運動量積分で排除厚さ δ* 補正を加算。",
    result="曲率連続コンタ + δ* 補正で出口一様流を実現するプログラム CONTUR を確立。",
    insight="出発点ワークフローの正体 = 世界標準。本提案は MOC を「ループ内形状生成器」へ昇格。中心線速度は一般 5 次多項式 (原典 Eq.35、端点 1・2 階微分整合で 3〜4 次に縮退)、スロート曲率は独立入力 RC (Mach6 で≈5.5) ＝原典精読で確認。",
    img="figs_final_sivells1978.png", img_cap="出典: Sivells AEDC-TR-78-63 (1978) Fig.1 — 変曲点で放射流となる Foelsch 型ノズル (C=変曲点, D=出口・一様流)")
add(sec="S1", short="標準教科書 — 回転流 MOC を含む体系", year="1977", status="○",
    cite="M. J. Zucrow & J. D. Hoffman, Gas Dynamics Vol. 2, Wiley, 1977 / J. D. Anderson, Modern Compressible Flow.",
    obj="圧縮性ガス力学と MOC を体系的に教科書化する。",
    method="最小長ノズル・軸対称 MOC・回転流 MOC (エントロピー勾配ソース項)・自由境界条件を定式化。",
    result="MOC 設計の標準リファレンスを提供。",
    insight="NS 由来 starting line ハイブリッドの回転流 MOC、SERN の自由境界処理の典拠。")

# ===== 第2章 スロート・収縮の物理 =====
add(sec="S2", short="スロート対流熱伝達の簡易相関", year="1957", status="○",
    cite='D. R. Bartz, "A Simple Equation for Rapid Estimation of Rocket Nozzle Convective Heat Transfer Coefficients," Jet Propulsion, 27:49–51, 1957.',
    obj="ロケットノズルの対流熱伝達係数を簡便に推算する。",
    method="境界層相関に基づく簡易式を提案。",
    result="スロート曲率は (D*/rc)^0.1 で効き hg ∝ rc^(-0.1)、小 rc ほど熱流束が増す。",
    insight="①の「熱流束 vs 一様性 vs 軸長」を rc で天秤にかける物理根拠。熱流束目的の一次見積りに有用。")
add(sec="S2", short="円錐超音速ノズルの熱伝達と再層流化", year="1967", status="◎",
    cite='L. H. Back, P. F. Massier, R. F. Cuffel, "Flow Phenomena and Convective Heat Transfer in a Conical Supersonic Nozzle," J. Spacecraft and Rockets, 4(8):1040–1047, 1967.',
    obj="円錐超音速ノズルの熱伝達機構を実験的に解明する。",
    method="壁熱流束・圧力を計測。",
    result="強い順圧力勾配下でスロート部の乱流境界層が再層流化し、壁熱流束が大幅に低下。",
    insight="スロート/膨張部の形状が熱負荷を左右する実験的裏付け。熱流束目的では遷移モデル化に留意。")
add(sec="S2", short="小スロート曲率の遷音速流れ場 (実験)", year="1969", status="◎",
    cite='R. F. Cuffel, L. H. Back, P. F. Massier, "Transonic Flowfield in a Supersonic Nozzle with Small Throat Radius of Curvature," AIAA Journal, 7(7):1364–1366, 1969.',
    obj="小スロート曲率 (rc/rth) の遷音速流れ場を実測する。",
    method="種々の rc/rth でスロート近傍の圧力・マッハ分布を計測。",
    result="rc/rth=0.625 で中心線静圧が壁の最大 3 倍、スロート面マッハが軸 0.8〜BL 端 1.4 と強い非一様。古典 2D 理論は rc/rth≳2 のみ妥当。",
    insight="小 rc 域で解析遷音速解が破綻 → NS 由来 starting line の必要根拠。rc をパレート変数化 (一様性側)。",
    img="figs_final_cuffel1969.png", img_cap="出典: Cuffel, Back & Massier (1969) AIAA J Fig.2 — rc/rth=0.625 スロート近傍の実験マッハ数コンタと理論 (Shelton/Prozan) の比較")
add(sec="S2", short="流量係数 Cd とスロート曲率", year="1975+", status="○",
    cite="A. J. Szaniszlo, NASA TN D-7848, 1975 / Alam ら 2016 / Li ら 2020 (ISO-9300 計量ノズル系).",
    obj="スロート曲率が流量係数 Cd に与える影響を定量化する。",
    method="計量ノズルの Cd 実測 (Szaniszlo: N2 100 atm, 実在気体補正, スロート Re 最大 8e6)。曲率影響は ISO-9300 系 (Alam, Li) で補完。",
    result="Szaniszlo は Cd を主にスロート Reynolds 数の関数として実測 (2 幾何で高 Re ほど漸増)。Cd が rc ~ 2–2.5 dth で最大・急曲率で剥離低下は ISO-9300 系 (Alam, Li) の曲率スイープ知見。",
    insight="Cd をパレートの一軸/制約にする根拠。限界: 計量ノズル系でロケット/極超とは間接的。",
    img="figs_final_szaniszlo1975.png", img_cap="出典: Szaniszlo (1975) NASA TN D-7848 Fig.4,5 — 2 幾何 (long-radius ASME / 連続壁曲率) の実験 Cd 対スロート Reynolds 数 (N2, 100 atm)")
add(sec="S2", short="軸対称収縮部の包括的設計法 (matched cubic)", year="1975", status="◎",
    cite='T. Morel, "Comprehensive Design of Axisymmetric Wind Tunnel Contractions," ASME J. Fluids Engineering, 97(2):225–233, 1975.',
    obj="剥離せず一様ソニックラインを得る軸対称収縮部を設計する。",
    method="matched cubic (2 本の 3 次曲線を変曲点で接続) を 1 パラメータ族とし非粘性ポテンシャル解析。",
    result="入口/出口の壁圧係数から剥離・出口非一様を制御する設計法を確立。",
    insight="亜音速収縮の標準初期形状。収縮比・長さを設計変数に開放する際の探索中心。")
add(sec="S2", short="5 次多項式収縮設計", year="1988", status="◎",
    cite="J. H. Bell & R. D. Mehta, Contraction Design for Small Low-Speed Wind Tunnels, NASA CR-177488, 1988.",
    obj="小型低速風洞の収縮部設計法を確立する。",
    method="5 次多項式壁 + 3D パネル法ポテンシャル場 + 2D 境界層コードで剥離予測の反復設計。",
    result="剥離なく低非一様な収縮コンタを得る。原典: 壁 Y=H_i-(H_i-H_e)[6ξ⁵-15ξ⁴+10ξ³]、7 次・マッチドキュービックは入口剥離 → 5 次採用。",
    insight="粘性比較で最も低非一様 → 収縮初期形状の既定候補。罠: 端点曲率 0 条件のためスロート円弧へ C2 ブレンド必須。歴史的「5 次」の置き場所は収縮壁 (本図) と中心線分布 (Sivells) の 2 箇所＝スロート壁ではない。",
    img="figs_final_bellmehta1988.png", img_cap="出典: Bell & Mehta NASA CR-177488 (1988) Fig.13 — 収縮壁形状の比較 (3/5/7 次多項式・対称マッチドキュービック)")
add(sec="S2", short="Morel 収縮法の実験検証", year="1988", status="◎",
    cite='E. G. Tulapurkara & V. V. K. Bhalla, "Experimental Investigation of Morel\'s Method for Wind Tunnel Contractions," ASME J. Fluids Engineering, 110(1):45–47, 1988.',
    obj="Morel 収縮法の妥当性を実験で検証する。",
    method="面積比 12 / 3.464 の収縮を製作し流れ場を計測。",
    result="剥離なし・薄い出口境界層・低い出口非一様を確認。",
    insight="古典収縮レシピの信頼性担保。初期値・探索範囲の中心として安全に使える。")
add(sec="S2", short="収縮 3 形状の head-to-head 比較", year="2017", status="◎",
    cite='Hassan/Zanoun ら, "Flow characteristics in low-speed wind tunnel contractions: Simulation and testing," Alexandria Engineering Journal, 2017.',
    obj="代表的な 3 つの収縮形状を統一条件で直接比較する。",
    method="5 次多項式・二重 3 次円弧 (Morel 型)・Witoszynski を CFD + LDA 実験で比較。",
    result="5 次多項式が最も低非一様 (<0.5%)、Witoszynski は壁近傍で圧力勾配が急変。",
    insight="収縮初期形状の既定は Bell–Mehta 5 次多項式。限界: 低速比較で圧縮性・高総温は要追確認。")
add(sec="S2", short="ICN とスロート非対称性の影響", year="2022", status="◎",
    cite='Rafiq, Rasheed, Afzal, Masoodi, "Influence of ideal nozzle geometry on supersonic flow using the method of characteristics," J. Mechanical Science and Technology, 36:6027–6039, 2022.',
    obj="理想コンタノズル (ICN) 幾何が超音速流に与える影響を解析する。",
    method="円弧スロート (Ru, Rd) + MOC 旋回コンタを構築・評価。",
    result="非対称スロート (Ru≠Rd) はソニックラインのずれから急マッハ変化・弱衝撃を生み流れ品質を劣化。",
    insight="スロート円弧の区分構成の典拠。Ru/Rd 非対称を変数化する際の品質リスクの示唆。")
add(sec="S2", short="スロート曲率と熱化学侵食 (最近接の rc 研究)", year="2015", status="◎",
    cite='D. Bianchi, F. Nasuti, M. Onofri, "Radius of Curvature Effects on Throat Thermochemical Erosion in Solid Rocket Motors," J. Spacecraft and Rockets, 52(2):320–330, 2015.',
    obj="スロート曲率半径 rc が熱化学侵食に与える効果を定量化する。",
    method="検証済み NS + 有限速度化学アブレーションで rc を振る (SRM)。",
    result="rc 低減が侵食低減・性能向上をもたらすことを定量化。",
    insight="rc 単一効果の最近接前例。だが多目的最適化でない → 「rc をパレート変数化」は文献空白 = 本ツールの新規性 (3-0 確証)。")

# ===== 第3章 NS-in-the-loop =====
add(sec="S3", short="NS-in-the-loop ノズル設計の直接前例 (CAN-DO)", year="1992", status="◎",
    cite='J. J. Korte, "Aerodynamic Design of Axisymmetric Hypersonic Wind-Tunnel Nozzles Using a Least-Squares/PNS Procedure," J. Spacecraft & Rockets, 29(5):685–691, 1992 / AIAA 92-4009 (CAN-DO).',
    obj="排除厚さ補正に頼らず、粘性込みでノズルコンタを直接設計する。",
    method="亜音速・遷音速スロートをフル NS、超音速を PNS で解き、コンタを非線形最小二乗で決定。",
    result="「粘性方程式に基づくため別途補正が不要」= 排除厚さ補正の原理的不要化を実証。",
    insight="本ツール構想の直接前例。限界: 単目的・NS/PNS 分割 → forge は単一ソルバ全域 + 多目的で超える。",
    img="figs_final_korte_pns.png", img_cap="出典: Korte (1990) NASA TP-3050 — explicit-upwind PNS による極超音速インレット内部流の圧力/マッハ場。本図は CAN-DO 本体でなく PNS アルゴリズム基盤")
add(sec="S3", short="超/極超風洞ノズルコンタ設計の現代運用", year="2006", status="◎",
    cite='W. Shope, "Contour Design Techniques for Super/Hypersonic Wind Tunnel Nozzles," AIAA 2006-3665, 2006.',
    obj="超/極超風洞ノズルコンタ設計技術を現代的に整理する。",
    method="AEDC における MOC ベース設計手法をレビュー。",
    result="Sivells 系手法が現役の世界標準であることを示す。",
    insight="ベースライン (低忠実度層) 手法の参照。")
add(sec="S3", short="収縮部の Bézier パラメータ化 + EGO 最適化", year="2007", status="◎",
    cite='C. Doolan & R. C. Morgans, "Numerical Evaluation and Optimization of Low Speed Wind Tunnel Contractions," AIAA 2007-3827, 2007.',
    obj="収縮部を固定式でなく最適化変数として設計する。",
    method="収縮壁を 2 パラメータ Bézier 化し、ポテンシャル流 + 積分境界層 (Thwaites) を SQP/DIRECT/EGO で最適化。",
    result="自動最適化で剥離なく良好な収縮を取得。",
    insight="収縮を最適化変数にする前例。限界: 低速・収縮部のみ → 本ツールは全域・圧縮性へ拡張。")
add(sec="S3", short="サロゲート支援 EA × 化学反応 RANS (scramjet)", year="2012", status="◎",
    cite='H. Ogawa & R. R. Boyce, "Nozzle Design Optimization for Axisymmetric Scramjets by Using Surrogate-Assisted Evolutionary Algorithms," J. Propulsion and Power, 28(6):1324–1338, 2012.',
    obj="軸対称スクラムジェットノズルを CFD ループ内で最適化する。",
    method="エリート実数値 GA (個体数 32, SBX+多項式突然変異) + kriging/RBF サロゲート (5 世代毎再学習) + 有限速度化学 RANS (25 反応/12 化学種, Menter 2 方程式 k-ω)。",
    result="燃料 on でベル形・off で円錐に近い最適形状、2 高度で評価。",
    insight="本提案フェーズ A (サロゲート EA × 粘性 RANS) の実証。温度依存・多成分を持つ最近接実装 (単目的推力)。",
    img="figs_final_ogawa2012.png", img_cap="出典: Ogawa & Boyce (2012) JPP Fig.19 — Baseline と Optimum 軸対称スクラムジェットノズルのマッハ数分布 (上: fuel-on / 下: fuel-off, 27 km)")
add(sec="S3", short="SERN の設計空間探索 (data mining)", year="2013", status="◎",
    cite='W. Huang ら, "Design exploration for a single expansion ramp nozzle (SERN) using data mining," Acta Astronautica, 83:10–17, 2013.',
    obj="SERN の多目的トレード構造を把握する。",
    method="data mining で推力 / 揚力 / ピッチモーメントの関係を探索。",
    result="各目的間のトレードオフ構造を可視化。",
    insight="③ SERN 多目的トレードの先行知見。3D 効果 (角部 R・サイドフェンス) は未踏 = 本ツールの拡張領域。")
add(sec="S3", short="CFD ベース境界層補正 (BL 端 95% 定義)", year="2014", status="◎",
    cite='"Boundary Layer Correction of Hypersonic Wind-tunnel Nozzle Designed by the Methods of Characteristics," J. Korean Soc. Aeronautical & Space Sciences, 42(12):1028–1036, 2014.',
    obj="MOC 非粘性コンタの境界層補正を粘性 CFD ベースで行う。",
    method="解析 δ* 式の代わりに粘性 CFD から BL 厚さを取得し、複数の BL 端定義を比較。",
    result="「断面最大速度の 95%」を BL 端とする定義が設計マッハ数回復に最良。",
    insight="「補正 → 確認」から NS ループ内在化への橋渡し。forge で δ* を直接取る設計 (B2.5) の典拠。")
add(sec="S3", short="超音速風洞ノズルの多目的最適化 (最重要近接例)", year="2022", status="◎",
    cite='K. Matsunaga, K. Fujio, H. Ogawa, H. Higa, K. Handa, "Nozzle design optimization for supersonic wind tunnel by using surrogate-assisted evolutionary algorithms," Aerospace Science and Technology, 130:107879, 2022.',
    obj="超音速風洞ノズルを多目的最適化する。",
    method="サロゲート支援 EA + CFD、目的 = マッハ偏差 vs 流れ偏向 (ノズル長とのトレード)。",
    result="一様性と偏向のパレート前線を取得。",
    insight="①の超音速版に直結 = 「風洞一様性のモダン最適化」前例の確証。限界: スロート曲率 rc は固定。")
add(sec="S3", short="B-spline による ED ノズル多目的最適化", year="2024", status="◎",
    cite='"Multi-objective aerodynamic optimization of expansion-deflection nozzle based on B-spline curves," Aerospace Science and Technology, 2024.',
    obj="膨張偏向 (ED) ノズルを多目的最適化する。",
    method="B-spline 制御点で壁を表現し、RANS + RBF サロゲート + NSGA-II。複数 NPR で推力効率を評価。",
    result="複数 NPR にわたるパレート解を取得。",
    insight="壁スプライン直接法 + NSGA-II の実装手本。限界: 超音速コンタのみ (収縮・スロートは固定)。")
add(sec="S3", short="POD + Kriging 場予測による極超推力ノズル最適化", year="2024", status="◎",
    cite='Huang, Wang, Wu, "A surrogate-based flow-field prediction and optimization strategy for hypersonic thrust nozzle," AIP Advances, 14:125312, 2024.',
    obj="極超音速推力ノズルの最適化を場予測で加速する。",
    method="NURBS コンタ + POD + Kriging で流れ場を予測するサロゲート最適化。",
    result="場サロゲートで評価コストを削減。",
    insight="場サロゲート活用の現代例。「サロゲートは補助、最終判定は CFD」運用と整合。")
add(sec="S3", short="ハイブリッドサロゲートによるロケットノズル多目的最適化", year="2025", status="◎",
    cite='Zhang ら, "A multi-objective optimization approach for rocket nozzle design based on hybrid surrogate model," Physics Letters A, 2025.',
    obj="ロケットノズルを多目的最適化する。",
    method="quintic + cubic-Bézier コンタ、MOEA/D で推力最大 + 長さ最小。",
    result="推力 vs 長さのパレート前線を取得。",
    insight="②「軸長 vs 推力」パレートの直接前例。rc・収縮は不変 = 全域統合は依然文献空白 (本ツール新規性)。")

# ===== 第4章 随伴法 =====
add(sec="S4", short="制御理論による空力設計 (随伴法の起点)", year="1988", status="◎",
    cite='A. Jameson, "Aerodynamic Design via Control Theory," J. Scientific Computing, 3:233–260, 1988.',
    obj="形状最適化を制御理論として定式化する。",
    method="随伴 PDE 1 回の求解で目的汎関数の形状勾配を取得 (非粘性)。",
    result="勾配計算コストが設計変数の数に依存しないことを示す。",
    insight="随伴法の核心 (変数数によらない勾配) の確立。",
    img="figs_final_jameson1988.png", img_cap="出典: Jameson (1988) Fig.1 — z 平面の翼型を σ 平面の単位円へ等角写像 (制御理論的形状最適化の定式化基盤)")
add(sec="S4", short="Navier-Stokes 方程式への随伴拡張", year="1998", status="◎",
    cite='A. Jameson, L. Martinelli, N. A. Pierce, "Optimum Aerodynamic Design Using the Navier-Stokes Equations," Theoretical and Computational Fluid Dynamics, 10:213–237, 1998.',
    obj="随伴法を粘性 NS へ拡張する。",
    method="NS 方程式に対する随伴方程式を導出・実装。",
    result="1 設計サイクル ≒ 流れ計算 2 回。粘性 (NS) 再設計を約 10 サイクルで達成 (Euler 再設計は約 60 サイクル)。",
    insight="「NS を回しながらの形状最適化」を可能にする核心 (MOC では原理的に不可)。フェーズ B の理論基盤。",
    img="figs_final_jameson1998.png", img_cap="出典: Jameson, Martinelli & Pierce (1998) Fig.5 — 広胴輸送機翼の粘性 (NS) 再設計。初期 vs 再設計形状と上面 Cp (Cd 0.0199→0.0194, M=0.83)")
add(sec="S4", short="SU2 — オープンソース随伴設計スイート", year="2016", status="◎",
    cite='T. D. Economon, F. Palacios, S. R. Copeland, T. W. Lukaczyk, J. J. Alonso, "SU2: An Open-Source Suite for Multiphysics Simulation and Design," AIAA Journal, 54(3):828–846, 2016.',
    obj="随伴設計を統合したオープンソース CFD スイートを提供する。",
    method="連続随伴・離散随伴 (AD)・Hicks-Henne/FFD・弾性メッシュ変形・SLSQP を同一コードに統合。",
    result="一気通貫の空力形状最適化フレームワークを実現。",
    insight="最有力の実装手本。自作前に「SU2 を勾配エンジンとして併用」する現実解 (フェーズ B-i)。",
    img="figs_final_su2_2016.png", img_cap="出典: Economon ら (2016) Fig.9 — DLR-F6 機体の Cp と随伴で算出した表面感度コンター (CL・CD・CMy 感度)")
add(sec="S4", short="AD による離散随伴 (CoDiPack)", year="2016", status="◎",
    cite='T. Albring, M. Sagebaum, N. R. Gauger, "Efficient Aerodynamic Design using the Discrete Adjoint Method in SU2," AIAA 2016-3518.',
    obj="フロー反復を逆微分して離散随伴を構築する。",
    method="リバースモード AD (式テンプレート CoDiPack) をフロー反復に適用。",
    result="t_adjoint / t_flow = 1.17 の実用コスト、モデル・目的関数の拡張が容易。",
    insight="forge に随伴を内製する場合 (B-ii) の推奨手法。多成分 TP・SST まで AD を通す設計の指針。",
    img="figs_final_albring2016.png", img_cap="出典: Albring, Sagebaum & Gauger (2016) Fig.3 — SU2 離散随伴のトップレベル構造。フロー 1 反復を逆微分し随伴を不動点反復")
add(sec="S4", short="極超音速インレットへの連続随伴適用 (SU2)", year="2017", status="◎",
    cite='Kline & Alonso, "Adjoint of Generalized Outflow-Based Functionals Applied to Hypersonic Inlet Design," AIAA Journal, 55(11):3903–3915, 2017.',
    obj="極超音速インレットを随伴で設計する。",
    method="出口面の一般化汎関数 (マッハ一様度等) に対する連続随伴。非粘性・粘性 RANS の両条件で勾配を検証。",
    result="高速流での随伴最適化を実証。",
    insight="高速流随伴の実用性の傍証。出口面汎関数 (マッハ一様度) の随伴定式化の直接の参考。",
    img="figs_final_kline_alonso.png", img_cap="出典: Kline & Alonso AIAA-2015-3060 Fig.13 — Mach 7 インレット最適化。(a) 初期 / (b) 最適化後のマッハ場 (衝撃をカウルリップへ整列) / (c) FFD 制御点")
add(sec="S4", short="研究コードへの離散随伴後付け実装 (Eilmer)", year="2020", status="◎",
    cite='K. Damm, R. J. Gollan, P. A. Jacobs, M. K. Smart, "Discrete Adjoint Optimization of a Hypersonic Inlet," AIAA Journal, 58(6):2621–2634, 2020.',
    obj="研究機関の内製コードに離散随伴を後付けする。",
    method="高速流 RANS の離散随伴を研究コード Eilmer に実装。",
    result="NASA P2 極超インレットの反射衝撃波を最適化でほぼ除去 (Bézier 形状制御)。",
    insight="「内製 CFD への離散随伴後付け」の直接前例 = forge の AD 化 (B-ii) の実現可能性の証拠。",
    img="figs_final_damm2020.png", img_cap="出典: Damm (2020) 博士論文 Fig.8.13 — NASA P2 極超インレットの圧力コンタ。(a) ベースライン / (b) 最適化後 (20 点 Bézier)、反射衝撃をほぼ除去")

# ===== 第5章 パラメータ化 =====
add(sec="S5", short="Hicks–Henne バンプ関数", year="1978", status="◎",
    cite='R. M. Hicks & P. A. Henne, "Wing Design by Numerical Optimization," J. Aircraft, 15(7):407–412, 1978.',
    obj="既存翼形状を数値最適化で改良する。",
    method="滑らかな局所バンプ関数の重ね合わせで形状を摂動。",
    result="少パラメータで局所形状を制御。",
    insight="随伴での多変数微調整 (既存コンタの polish) に適合。SU2 の 2D 既定パラメータ化。")
add(sec="S5", short="FFD — 自由形状変形", year="1986", status="○",
    cite='T. W. Sederberg & S. R. Parry, "Free-Form Deformation of Solid Geometric Models," SIGGRAPH, 20(4):151–160, 1986.',
    obj="任意形状を直感的に変形する手法を確立する。",
    method="形状を包む制御格子の変形で内部形状を連続変形。",
    result="CAD フリーで任意トポロジに適用可能。",
    insight="③ SERN の 3D・角部 R に最適 (forge メッシュを箱で包んで変形)。随伴と直結。")
add(sec="S5", short="CST — Class-Shape Transformation", year="2008", status="○",
    cite='B. M. Kulfan, "Universal Parametric Geometry Representation Method," J. Aircraft, 45(1):142–158, 2008.',
    obj="翼/断面形状を少変数で普遍的に表現する。",
    method="クラス関数 × シェイプ関数で形状族を構成。",
    result="丸み先端 (class 指数 N1=0.5) で前縁特異点を解析的に表現可能。",
    insight="①② 軸対称コンタの代替パラメータ化 (少変数 → 進化計算が回る)。",
    img="figs_final_kulfan2008.png", img_cap="出典: Kulfan AIAA-2007-62 (2008) Fig.6 — 翼断面の単位シェイプ関数を Nose Shape (S1) と Aft-end Shape (S2) に分解 (class×shape)")
add(sec="S5", short="CAD フリー高忠実度最適化 (pyGeo FFD)", year="2010", status="◎",
    cite='G. K. W. Kenway & J. R. R. A. Martins, "A CAD-Free Approach to High-Fidelity Aerostructural Optimization," AIAA 2010-9231.',
    obj="CAD フリーで高忠実度の空力構造最適化を行う。",
    method="FFD ベースの形状制御 (pyGeo) + 随伴。",
    result="大規模設計変数の最適化を実証。",
    insight="「FFD + 随伴 + メッシュワープ」標準分業 (MACH-Aero) の形状層の手本。",
    img="figs_final_kenway2010.png", img_cap="出典: Kenway & Martins AIAA-2010-9231 (2010) — FFD 制御箱で機体全体を包み、表面メッシュ点 (赤) を制御点で連続変形")

# ===== 第6章 サロゲート・多目的 =====
add(sec="S6", short="EGO — ベイズ最適化の原点", year="1998", status="○",
    cite='D. R. Jones, M. Schonlau, W. J. Welch, "Efficient Global Optimization of Expensive Black-Box Functions," J. Global Optimization, 13:455–492, 1998.',
    obj="高コストな黒箱関数を少ない評価回数で大域最適化する。",
    method="Kriging + Expected Improvement で次の評価点を選定。",
    result="少ない評価で大域最適へ収束。",
    insight="expensive な forge 評価と最も相性の良い infill 戦略の基礎。",
    img="figs_final_jones_ego1998.png", img_cap="出典: Jones ら (1998) Fig.12 — Branin 関数の等高線。初期標本 (白) と期待改善最大化で選ばれた最初の 3 点 (黒)")
add(sec="S6", short="多忠実度 co-kriging の起点", year="2000", status="○",
    cite='M. C. Kennedy & A. O\'Hagan, "Predicting the output from a complex computer code when fast approximations are available," Biometrika, 87(1):1–13, 2000.',
    obj="高速近似と少数の高忠実度評価を統計的に融合する。",
    method="自己回帰モデル z_t=ρ·z_(t-1)+δ_t でレベル間をガウス過程結合。",
    result="低コストで高忠実度予測を実現。",
    insight="MOC (秒) / 軸対称 RANS (分) / 3D RANS (時間) の多忠実度構成 (B3) の理論基盤。",
    img="figs_final_kennedy2000.png", img_cap="出典: Kennedy & O'Hagan (2000) Fig.2 — 3 忠実度レベルの模擬コード出力 z(x,t1)/z(x,t2)/z(x,t3)")
add(sec="S6", short="NSGA-II / NSGA-III — 多目的進化計算の標準", year="2002/14", status="○",
    cite='K. Deb ら, "A Fast and Elitist Multiobjective Genetic Algorithm: NSGA-II," IEEE TEC, 6(2):182–197, 2002 / "NSGA-III," IEEE TEC, 18(4):577–601, 2014.',
    obj="多目的最適化でパレート前線全体を一度の探索で得る。",
    method="非劣ソート + 混雑度 (III は参照点) で進化計算。",
    result="非凸前線・形状制約に強いパレート探索を実現。",
    insight="パレート要件の主役 (フェーズ A)。限界: 評価数が数千 → 粘性 NS 直結は高コスト → サロゲート前提。",
    img="figs_final_nsga2002.png", img_cap="出典: Deb ら (2002) Fig.2 — NSGA-II 手順: 非劣ソートで前線 F1,F2,F3 に層別し混雑度ソートで次世代を選抜")
add(sec="S6", short="多目的ベイズ最適化 (ParEGO / EHVI / MOEA/D)", year="2006+", status="○",
    cite='J. Knowles, "ParEGO," IEEE TEC, 10(1), 2006 / Emmerich ら (EHVI) / Zhang & Li, "MOEA/D," IEEE TEC, 11(6), 2007.',
    obj="expensive な多目的最適化を効率化する。",
    method="重み付き Tchebycheff スカラー化+EGO (ParEGO) / 超体積期待改善 (EHVI) / 分解 (MOEA/D)。",
    result="少評価で多目的パレートを探索。",
    insight="フェーズ A の infill 戦略候補群。各手法の直接比較は未検証 → 採用前に追検証。",
    img="figs_final_parego2006.png", img_cap="出典: Knowles (2006) Fig.3 — 目的空間 (f1 vs f2) の非劣点 5 集合と best/median/worst 到達面")
add(sec="S6", short="サロゲートによる多忠実度最適化", year="2007/08", status="◎",
    cite='A. I. J. Forrester, A. Sóbester, A. J. Keane, "Multi-Fidelity Optimization via Surrogate Modelling," Proc. R. Soc. A, 463:3251–3269, 2007.',
    obj="複数忠実度の評価をサロゲートで融合し最適化コストを削減する。",
    method="co-kriging で低忠実度多数 + 高忠実度少数を結合。",
    result="安価に大域近似を構築。",
    insight="MOC を低忠実度層に使う本提案 B3 (多忠実度予算配分) の直接の方法論。",
    img="figs_final_forrester2007.png", img_cap="出典: Forrester ら (2007) Fig.1 — 一変数 co-kriging 例。高忠実度 fe・低忠実度 fc・素の Kriging・co-kriging 予測の比較")
add(sec="S6", short="GEK — 勾配強化 Kriging", year="2008", status="◎",
    cite='J. Laurenceau & P. Sagaut, "Building Efficient Response Surfaces of Aerodynamic Functions with Kriging and Cokriging," AIAA Journal, 46(2):498–507, 2008.',
    obj="勾配情報を使って応答曲面の精度を上げる。",
    method="勾配ベクトルを補間に取り込む gradient-enhanced Kriging (GEK)。",
    result="素の Kriging より精度を劇的に改善 (2–6 次元の変形遷音速翼で実証)。",
    insight="随伴勾配 (単目的) を多目的パレート探索に活かす機構の実証 (その 1)。")
add(sec="S6", short="Weighted GEK — 高次元への拡張", year="2017", status="◎",
    cite='Z.-H. Han, Y. Zhang, C.-X. Song, K.-S. Zhang, "Weighted Gradient-Enhanced Kriging for High-Dimensional Surrogate Modeling and Design Optimization," AIAA Journal, 55(12):4330–4346, 2017.',
    obj="GEK を高次元設計に拡張する。",
    method="随伴勾配を Kriging に注入、小モデル分割 + 重み付き和で次元の呪いを緩和。",
    result="36–108 変数の RANS + 随伴の遷音速翼逆設計を実証。",
    insight="フェーズ B (随伴) とフェーズ A (サロゲート MOO) を繋ぐ実証済みの機構 (その 2)。")
add(sec="S6", short="深層学習サロゲート (DeepONet / MeshGraphNets)", year="2020/21", status="△",
    cite="Lu, Jin, Karniadakis ら, DeepONet, Nature Machine Intelligence, 2021 / Pfaff ら, MeshGraphNets, 2020.",
    obj="流れ場を ML (作用素学習 / GNN) で高速予測する。",
    method="DeepONet (作用素学習)・MeshGraphNets (GNN)。",
    result="数値ソルバより 1–2 桁高速だが信頼性・外挿は未成熟。",
    insight="「サロゲートは infill 補助に限定、最終判定は forge」運用が安全。設計成熟度は未検証。")

# ===== 第7章 参照アーキテクチャ =====
add(sec="S7", short="MACH-Aero — 随伴最適化の標準分業", year="–", status="◎",
    cite="MACH-Aero (U. Michigan MDO Lab) 公式ドキュメント: pyHyp / pyGeo / IDWarp / ADflow / DAFoam / pyOptSparse.",
    obj="随伴最適化の標準ツールチェーンを提供する。",
    method="6 モジュール分業 (前処理/FFD 形状/メッシュ変形/流れ + 随伴/最適化)。",
    result="大規模空力構造最適化の実績。",
    insight="「随伴 + FFD + メッシュワープ + SQP」標準分業の手本。forge ラッパ層のモジュール設計指針。")
add(sec="S7", short="OpenMDAO — 多分野最適化フレームワーク", year="2019", status="◎",
    cite='J. S. Gray, J. T. Hwang, J. R. R. A. Martins, K. T. Moore, B. A. Naylor, "OpenMDAO: an open-source framework for MDAO," Structural and Multidisciplinary Optimization, 59(4), 2019.',
    obj="多分野連成最適化を効率化する OSS 基盤を提供する。",
    method="統一微分方程式 (MAUD) で全分野を単一陰系として解き順方向/随伴を統一算出。",
    result="MDO の統一フレームワークを実現。",
    insight="①凝縮/熱流束・②多高度/バルブ過渡・③3D/運転条件の多目的を束ねる基盤の第一候補。",
    img="figs_final_openmdao2019.png", img_cap="出典: Gray ら (2019) Fig.5 — OpenMDAO 自動生成のモデル可視化。左=階層ツリー、右=N² 依存グラフ")
add(sec="S7", short="DAKOTA — 汎用最適化・UQ ドライバ", year="–", status="○",
    cite="DAKOTA (Sandia National Laboratories). DOE・サロゲート・多目的 EA・UQ を統合した汎用基盤。",
    obj="汎用の最適化・不確かさ定量化 (UQ) ドライバを提供する。",
    method="DOE・サロゲート・多目的 EA・UQ を統合し、外部ソルバと疎結合。",
    result="SU2 等との連携実績多数。",
    insight="forge を黒箱として最短で回す選択肢 (OpenMDAO との二択)。")

# ===== 第8章 凝縮 =====
add(sec="S8", short="風洞凝縮の古典総説", year="1958", status="○",
    cite='P. P. Wegener & L. M. Mack, "Condensation in Supersonic and Hypersonic Wind Tunnels," Advances in Applied Mechanics, 5:307–447, 1958.',
    obj="超音速・極超音速風洞の凝縮現象を包括的に整理する。",
    method="風洞凝縮の理論・実験を総説。",
    result="凝縮の物理と試験への影響を体系化。",
    insight="①固有の凝縮問題の背景知識の標準参照。")
add(sec="S8", short="非平衡凝縮と凝縮ショックの理論", year="1976", status="◎",
    cite='J. P. Sislian & I. I. Glass, AIAA J., 14(12):1731–1737, 1976 / P. A. Blythe & C. J. Shih, "Condensation shocks in nozzle flows," J. Fluid Mech., 76(3):593–621, 1976.',
    obj="非平衡凝縮と凝縮ショックの形成を理論化する。",
    method="均質核生成理論 + 特性線解析。",
    result="急膨張 → 過飽和 → 核生成 → 潜熱放出 → 同族特性線交差で凝縮ショックが形成。",
    insight="凝縮が試験コアのマッハ一様性を乱す機構の理論基盤。")
add(sec="S8", short="空気・窒素の凝縮 onset 相関", year="1968", status="◎",
    cite='F. L. Daum & G. Gyarmathy, "Condensation of Air and Nitrogen in Hypersonic Wind Tunnels," AIAA Journal, 6(3):458–465, 1968.',
    obj="空気・窒素の凝縮 onset 条件を定量化する。",
    method="極超音速風洞で凝縮 onset を計測・相関。",
    result="低圧で空気は実質純 N2、「局所膨張率 / 静圧」比で onset を相関できる。",
    insight="dMc/dx への上限制約として Bézier 分布上で CFD 前にチェック可 (B2.5 事前フィルタ)。凝縮余裕制約の根拠。")
add(sec="S8", short="極超音速ノズルの N2 凝縮 CFD", year="2014", status="◎",
    cite='Lin, Cheng, Luo ら, "On nitrogen condensation in hypersonic nozzle flows," Shock Waves, 24:179–189, 2014.',
    obj="極超音速ノズルの N2 凝縮を数値予測する。",
    method="古典核生成理論 (CNT) + Gyarmathy 液滴成長を 2D・軸対称 CFD に結合。",
    result="総温 T0 を上げると凝縮量・出口圧温変化が減少 = 凝縮余裕が増す。",
    insight="forge に凝縮モデルを載せる際の実装参考。総温を設計変数に含める根拠。",
    img="figs_final_lin2014.png", img_cap="出典: Lin ら (2014) Fig.3 — N2 凝縮の核生成率・液相質量分率コンター")
add(sec="S8", short="N2 均質凝縮の半経験モデル", year="2020", status="◎",
    cite='K. Lax & S. B. Leonov, "Homogeneous Condensation of Nitrogen in Hypersonic Wind Tunnels: A Semi-Empirical Model," AIAA Journal, 58(11):4807–4818, 2020.',
    obj="N2 均質凝縮の高精度 onset モデルを構築する。",
    method="DFT ベースの半経験モデル。",
    result="Mach 10 軸対称ノズルで onset 静温を 3.5 K 以内で再現。",
    insight="凝縮余裕制約の定量精度の現代水準。後処理ベース制約の較正に使用可能。")
add(sec="S8", short="凝縮指標を目的関数にした形状最適化", year="2017", status="◎",
    cite='M. Noori Rahim Abadi, A. Ahmadpour, J. P. Meyer ら, "CFD-based shape optimization of steam turbine blade cascade in transonic two phase flows," Applied Thermal Engineering, 112:1575–1589, 2017.',
    obj="凝縮指標を最適化目的に据えて形状最適化する。",
    method="核生成率・最大液滴半径を目的関数に蒸気タービン翼列を最適化。",
    result="効率 +2.1% を達成。",
    insight="「凝縮 onset を目的/制約にする」実証 (湿り蒸気限定)。極超風洞試験ガスは前例なし = 本ツールの新規貢献余地。")

# ===== 第9章 DACS・SERN =====
add(sec="S9", short="横噴流制御ミサイルの衝撃干渉", year="2006", status="◎",
    cite='B.-Y. Min, J.-W. Lee, Y.-H. Byun, "Numerical investigation of the shock interaction effect on the lateral jet controlled missile," Aerospace Science and Technology, 10(5):385–393, 2006.',
    obj="横噴流制御ミサイルの衝撃干渉を解明する。",
    method="3D RANS で噴流圧力比・マッハ数を変えて解析。",
    result="分離/弓状/バレル衝撃・マッハディスクの干渉構造、誘起法線力は推力に比例だが干渉で損失。",
    insight="②外部クロスフロー干渉評価シナリオの基礎。側力・モーメントを目的関数化する根拠。",
    img="figs_final_min2006.png", img_cap="出典: Min ら (2006) — 噴流マッハ数ごとの衝撃干渉構造 (マッハ数コンター)")
add(sec="S9", short="ピントルノズルの動特性 (可変推力)", year="2015", status="◎",
    cite='J. Heo, K. Jeong, H.-G. Sung, "Numerical Study of the Dynamic Characteristics of Pintle Nozzles for Variable Thrust," J. Propulsion and Power, 31(1):230–237, 2015.',
    obj="可変推力ピントルノズルの動特性を解明する。",
    method="スライディングメッシュ非定常 RANS。",
    result="往復・挿入・引抜時の推力応答遅れを定量化、リップ衝撃・剥離を捕捉。",
    insight="②バルブ開閉過渡を NS で解く前例。forge のデュアルタイム + 動的メッシュ要件の参考。",
    img="figs_final_heo2015.png", img_cap="出典: Heo ら (2015) — ピントルストロークごとの密度勾配 (シュリーレン風)")
add(sec="S9", short="ピントル弁付き SRM の過渡流れ", year="2020", status="◎",
    cite='A. Song ら, "Transient flow characteristics and performance of a solid rocket motor with a pintle valve," Chinese J. Aeronautics, 33:3189–3205, 2020.',
    obj="ピントル弁付き SRM の過渡流れと性能を解明する。",
    method="動的メッシュ RANS、冷走試験と比較。",
    result="ピントル運動で圧縮波・反射衝撃が発生、冷走試験との誤差 < 2%。",
    insight="viscous-NS-in-the-loop の過渡評価が実用精度に達していることの証拠。",
    img="figs_final_song2020.png", img_cap="出典: Song ら (2020) — ピントル位置ごとのマッハ数コンターと推力履歴")
add(sec="S9", short="ピントルスラスタ DACS のシステムモデル", year="2020", status="◎",
    cite='M. Ji & H. Chang, "Modeling and dynamic characteristics analysis on solid attitude control motor using pintle thrusters," Aerospace Science and Technology, 106, 2020.',
    obj="ピントルスラスタ DACS のシステム動特性をモデル化する。",
    method="開/閉過渡を start-up・pintle-moving・thrust-establishment の 3 段に分解。",
    result="応答時間を再現するシステムモデルを構築。",
    insight="②のシステムレベル要求 (応答時間) と CFD 評価の接続点。",
    img="figs_final_ji2020.png", img_cap="出典: Ji & Chang (2020) — DACS 動特性のサブシステム・ブロック線図")
add(sec="S9", short="Fluidic ピントルと実機 DACS ハードウェア", year="2024/19", status="◎",
    cite='Yan ら, Aerospace (MDPI), 11(3):243, 2024 / K. W. Naumann (Bayern-Chemie), "Hot Gas Nozzle-Valve Assembly for Continuously Operating DACS," AIAA 2019-3879.',
    obj="Fluidic ピントルと実機 DACS ハードウェアを評価する。",
    method="k-ω SST 動的メッシュ (Yan) / 実機開発・制御 (Naumann)。",
    result="4 過渡モード解析 / 連続作動 DACS 用ホットガスバルブを実証。",
    insight="②の設計要求 (連続作動・応答性) の実機水準の把握。")
add(sec="S9", short="スカーフドノズルのプルーム特性", year="2018", status="◎",
    cite='P. Behrouzi, J. J. McGuirk, C. Avenell, "Effect of Scarfing on Rectangular Nozzle Supersonic Jet Plume Flow Characteristics," AIAA Journal, 56(1):301–315, 2018.',
    obj="スカーフ (切り欠き) が矩形ノズルプルームに与える影響を解明する。",
    method="LDA 計測 + RANS で膨張条件を変えて比較。",
    result="切り欠きが出口静圧を変え縦渦二次流れを誘起、過膨張で四つ葉・不足膨張で矩形 + 分岐。",
    insight="②スカーフ角の変数化と側力評価の基礎。注意: RANS (S-A) は二次流れ速度を過大評価。",
    img="figs_final_behrouzi2018.png", img_cap="出典: Behrouzi ら (2018) — プルーム断面の LDA vs CFD (四つ葉状二次流れ)")
add(sec="S9", short="超音速横風中の音速噴流 (パラメトリック)", year="2022", status="◎",
    cite='C.-L. Qiao ら, "Parametric study on the sonic transverse jet in supersonic crossflow," Aerospace Science and Technology, 123:107472, 2022.',
    obj="超音速横風中の音速噴流をパラメトリックに解明する。",
    method="RANS で噴流圧力比 PR・噴射角 AoI を変える。",
    result="PR・AoI が支配、高 PR で衝撃干渉が不安定化、RCS 増幅率が 1 を切る場合あり。",
    insight="②多高度 = 過/不足膨張 PR スイープで増幅率・側力をパレート評価する設計シナリオの根拠。")

# ===== 第10章 ④デュアルベル (5機種版・追補) =====
add(sec="S10", short="デュアルベルの原理 — 高度補償ノズル", year="1994/2016", status="○",
    cite='M. Horn & S. Fisher (Rocketdyne), "Dual-Bell Altitude Compensating Nozzles," NASA NTRS 19940018584 / "Altitude Compensating Nozzle," NASA NTRS 20160008031.',
    obj="単一の固定ノズルで広い高度域にわたり推力効率を確保する (高度補償)。",
    method="内側コンタを 2 つの重ねたベル (基本ベル + 延長ベル) で構成し、変曲点 (contour inflection) で接続する。可動部なし。",
    result="低高度 = モード 1 (基本ベルのみ、変曲点で安定剥離)、高高度 = モード 2 (延長部まで全付着) と作動が切り替わり、各高度に適合した膨張比を得る。",
    insight="計画の④機種。基盤は③ベル、これに変曲位置・延長長さ・延長部壁圧分布を設計変数として付加 (B5 ④)。延長部は壁圧帰還で設計 (B3)。",
    img="figs_final_dualbell_geom.png", img_cap="出典: NASA TM-2013-216590 — デュアルベル 3D 形状 (First bell = 基本ベル / Second bell = 延長ベル)")
add(sec="S10", short="2 モード作動とモード遷移", year="2013", status="○",
    cite='NASA Dryden, "Proposed Flight Research of a Dual-Bell Rocket Nozzle," NASA TM-2013-216590 / NTRS 20140003943.',
    obj="デュアルベルの 2 モード作動とモード遷移を CFD/飛行で実証する研究を提案する。",
    method="軸対称 CFD でモード 1 (海面) / モード 2 (高空) のマッハ場と遷移条件 (設計遷移高度) を解析。",
    result="設計遷移高度で剥離点が変曲点から延長端へジャンプし、推力係数が段階的に上昇 (高度補償)。",
    insight="forge では軸対称 RANS × 2 作動点 (NPR) で評価 (B5 ④、目的 = 海面推力 vs 高空推力 vs 軸長)。遷移の非定常・横荷重は URANS 最終確認 (付録10)。",
    img="figs_final_dualbell_modes.png", img_cap="出典: NASA TM-2013-216590 — モード 1 (海面・変曲点剥離) / モード 2 (高空・全付着) のマッハ数コンター")
add(sec="S10", short="高度補償ノズルの中での位置づけと forge 適用", year="–", status="○",
    cite='高度補償ノズル形式の比較 (NASA TM-2013-216590 Fig.) / 本調査 B5 ④・付録10。',
    obj="デュアルベルを他の高度補償ノズル (エアロスパイク/プラグ・E-D・dual-expander) の中に位置づける。",
    method="形式比較。デュアルベルは「ベル 2 段」で機構が単純・既存ベル製造法を流用でき実装容易。",
    result="固定形状・可動部なしで高度補償を実現。エアロスパイク等に比べ設計・製造リスクが低い。",
    insight="計画では③ベルの帰還エンジン (B3) をそのまま流用し、変曲位置・延長壁圧分布を追加変数化するだけで④に拡張できる。残課題: 遷移ヒステリシス幅・横荷重の定量 (付録10、要 URANS)。",
    img="figs_final_dualbell_types.png", img_cap="出典: NASA TM-2013-216590 — 各種高度補償/先進ノズル形式の比較 (デュアルベル・E-D・プラグ等)")

PAPERS = P

# ---------------------------------------------------------------- ヘルパ
def png_size(path):
    with open(path, "rb") as f:
        head = f.read(26)
    if head[:8] != b"\x89PNG\r\n\x1a\n":
        return None
    w, h = struct.unpack(">II", head[16:24])
    return w, h

def set_jp(run, size=None, bold=None, italic=None, color=None):
    f = run.font
    f.name = JP_FONT
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", JP_FONT)
    if size is not None: f.size = Pt(size)
    if bold is not None: f.bold = bold
    if italic is not None: f.italic = italic
    if color is not None: f.color.rgb = color

def add_box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    return tb, tf

def add_rect(slide, x, y, w, h, fill, line=None, round_=False, line_w=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w: shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def footer(slide, sec_title, idx, total):
    _, tf = add_box(slide, 0.45, 7.12, 10.8, 0.3)
    r = tf.paragraphs[0].add_run()
    r.text = "超音速・極超音速ノズル設計技術 文献調査  |  " + sec_title
    set_jp(r, size=9, color=GRAY)
    _, tf2 = add_box(slide, 12.1, 7.12, 0.85, 0.3)
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
    r2 = tf2.paragraphs[0].add_run()
    r2.text = f"{idx} / {total}"
    set_jp(r2, size=9, color=GRAY)

def status_badge(slide, status, x=11.45, y=0.16):
    label, col = STATUS[status]
    badge = add_rect(slide, x, y, 1.5, 0.4, col, round_=True)
    tf = badge.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    set_jp(r, size=10.5, bold=True, color=WHITE)

# ---------------------------------------------------------------- 構築
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

plan = [("title", None), ("toc", None)]
for key in SECTIONS:
    plan.append(("divider", key))
    for p in PAPERS:
        if p["sec"] == key:
            plan.append(("paper", p))
plan.append(("gap", None))
plan.append(("proposal", None))

TOTAL = len(plan)
num = 0

for kind, payload in plan:
    num += 1
    slide = prs.slides.add_slide(BLANK)

    if kind == "title":
        add_rect(slide, 0, 0, 13.333, 7.5, RGBColor(0x10, 0x2A, 0x43))
        add_rect(slide, 0.0, 4.55, 13.333, 0.06, RGBColor(0xD9, 0x8E, 0x32))
        _, tf = add_box(slide, 0.9, 1.6, 11.5, 2.0)
        r = tf.paragraphs[0].add_run()
        r.text = "超音速・極超音速ノズル設計技術 文献調査"
        set_jp(r, size=38, bold=True, color=WHITE)
        p2 = tf.add_paragraph(); p2.space_before = Pt(14)
        r2 = p2.add_run()
        r2.text = "古典 MOC から NS-in-the-loop・随伴法・サロゲート多目的最適化まで"
        set_jp(r2, size=19, color=RGBColor(0xCF, 0xDC, 0xEA))
        p3 = tf.add_paragraph(); p3.space_before = Pt(4)
        r3 = p3.add_run()
        r3.text = "— 1 論文 1 スライド / 目的・方法・結果・得られた知見で整理 —"
        set_jp(r3, size=15, color=RGBColor(0x9F, 0xB4, 0xC8))
        _, tf3 = add_box(slide, 0.9, 4.9, 11.5, 1.9)
        for i, t in enumerate([
            "対象 (5 機種): ①風洞(軸対称)  ②風洞(矩形)  ③ベルスラスタ  ④デュアルベル  ⑤SERN",
            "目的: 粘性 CFD をループ内で回し多目的パレート解を得る次世代ノズル設計ツールの研究提案",
            "2026-06-14   |   出典: forge リポジトリ調査文書 (敵対的検証 25/25 主張確証)",
        ]):
            p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
            p.space_after = Pt(8)
            r = p.add_run(); r.text = t
            set_jp(r, size=14, color=RGBColor(0xCF, 0xDC, 0xEA))

    elif kind == "toc":
        add_rect(slide, 0, 0, 13.333, 0.9, RGBColor(0x10, 0x2A, 0x43))
        _, tf = add_box(slide, 0.45, 0.16, 11.0, 0.6)
        r = tf.paragraphs[0].add_run(); r.text = "構成と読み方"
        set_jp(r, size=22, bold=True, color=WHITE)
        _, tfl = add_box(slide, 0.55, 1.2, 7.7, 5.7)
        counts = {k: sum(1 for p in PAPERS if p["sec"] == k) for k in SECTIONS}
        for i, (k, (title, col)) in enumerate(SECTIONS.items()):
            p = tfl.paragraphs[0] if i == 0 else tfl.add_paragraph()
            p.space_after = Pt(9)
            r = p.add_run(); r.text = title
            set_jp(r, size=14.5, bold=True, color=col)
            r2 = p.add_run(); r2.text = f"   ({counts[k]} 件)"
            set_jp(r2, size=12, color=GRAY)
        p = tfl.add_paragraph(); p.space_before = Pt(6)
        r = p.add_run(); r.text = "まとめ  技術空白 (新規性) と提案アーキテクチャ"
        set_jp(r, size=14.5, bold=True, color=DARK)
        # 右: 読み方 + 凡例
        add_rect(slide, 8.55, 1.2, 4.3, 2.55, RGBColor(0xF2, 0xF4, 0xF7))
        _, tfr = add_box(slide, 8.78, 1.38, 3.9, 2.3)
        r = tfr.paragraphs[0].add_run(); r.text = "各スライドの構成"
        set_jp(r, size=13, bold=True, color=DARK)
        for lab, desc, col in [
            ("目的", "なぜその研究をしたか", FIELD_DEF[0][1]),
            ("方法", "何をどう解いたか", FIELD_DEF[1][1]),
            ("結果", "何が分かったか", FIELD_DEF[2][1]),
            ("知見", "本ツールへの含意・限界", FIELD_DEF[3][1]),
        ]:
            p = tfr.add_paragraph(); p.space_before = Pt(6)
            r = p.add_run(); r.text = "■ " + lab + "  "
            set_jp(r, size=12, bold=True, color=col)
            r2 = p.add_run(); r2.text = "— " + desc
            set_jp(r2, size=12, color=DARK)
        add_rect(slide, 8.55, 3.95, 4.3, 2.5, RGBColor(0xF2, 0xF4, 0xF7))
        _, tfs = add_box(slide, 8.78, 4.13, 3.9, 2.3)
        r = tfs.paragraphs[0].add_run(); r.text = "凡例 (引用の検証ステータス)"
        set_jp(r, size=13, bold=True, color=DARK)
        for lab, desc in [
            ("◎ 検証済", "敵対的検証 (3 票方式) で 3-0 確証"),
            ("○ 古典・標準", "分野の古典・標準文献 (高確信)"),
            ("△ 要追検証", "取得済だが独立検証は未実施"),
        ]:
            p = tfs.add_paragraph(); p.space_before = Pt(7)
            r = p.add_run(); r.text = lab + "  — "
            set_jp(r, size=11.5, bold=True, color=STATUS[lab[0]][1])
            r2 = p.add_run(); r2.text = desc
            set_jp(r2, size=11.5, color=DARK)
        _, tfn = add_box(slide, 8.78, 6.55, 3.9, 0.5)
        r = tfn.paragraphs[0].add_run()
        r.text = "★ papers/ に PDF がある論文は代表図を掲載。"
        set_jp(r, size=11, italic=True, color=GRAY)

    elif kind == "divider":
        title, col = SECTIONS[payload]
        add_rect(slide, 0, 0, 13.333, 7.5, col)
        _, tf = add_box(slide, 0.9, 2.7, 11.5, 1.4)
        r = tf.paragraphs[0].add_run(); r.text = title
        set_jp(r, size=31, bold=True, color=WHITE)
        names = [(p2["short"], p2["year"], "img" in p2) for p2 in PAPERS if p2["sec"] == payload]
        _, tf2 = add_box(slide, 0.95, 4.15, 11.6, 2.9)
        for i, (nm, yr, has_img) in enumerate(names):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p.space_after = Pt(3)
            r = p.add_run(); r.text = f"・{nm} ({yr})"
            set_jp(r, size=13, color=RGBColor(0xE8, 0xEC, 0xF2))
            if has_img:
                r2 = p.add_run(); r2.text = "  ★図"
                set_jp(r2, size=12, bold=True, color=RGBColor(0xFF, 0xD7, 0x6A))
        footer(slide, title, num, TOTAL)

    elif kind == "paper":
        d = payload
        sec_title, col = SECTIONS[d["sec"]]
        has_img = "img" in d and os.path.exists(os.path.join(HERE, d["img"]))
        # タイトル帯
        add_rect(slide, 0, 0, 13.333, 0.8, col)
        _, tf = add_box(slide, 0.45, 0.12, 9.4, 0.58)
        r = tf.paragraphs[0].add_run(); r.text = d["short"]
        set_jp(r, size=18, bold=True, color=WHITE)
        _, tfy = add_box(slide, 9.9, 0.2, 1.3, 0.45)
        tfy.paragraphs[0].alignment = PP_ALIGN.RIGHT
        ry = tfy.paragraphs[0].add_run(); ry.text = d["year"]
        set_jp(ry, size=14, bold=True, color=RGBColor(0xDD, 0xE5, 0xEE))
        status_badge(slide, d["status"])
        # 書誌
        _, tfc = add_box(slide, 0.45, 0.92, 12.45, 0.66)
        rc = tfc.paragraphs[0].add_run(); rc.text = d["cite"]
        set_jp(rc, size=10.5, italic=True, color=GRAY)

        # テキスト領域の幅 (画像があれば左側に寄せる)
        if has_img:
            text_w = 7.05
            label_x = 0.45
            text_x = 1.45
            txt_inner_w = text_w - (text_x - label_x)
        else:
            label_x = 0.45
            text_x = 1.55
            txt_inner_w = 11.35

        # 4 区分を縦に並べる
        y0 = 1.72
        pitch = 1.30 if not has_img else 1.28
        for fi, key in enumerate(FIELD_KEYS):
            ry_ = y0 + fi * pitch
            lab_text, lab_col = FIELD_DEF[fi]
            # ラベルタブ
            tab = add_rect(slide, label_x, ry_ + 0.02, 0.92, 0.44, lab_col, round_=True)
            ttf = tab.text_frame
            ttf.word_wrap = False
            ttf.margin_left = Emu(0); ttf.margin_right = Emu(0)
            ttf.margin_top = Emu(0); ttf.margin_bottom = Emu(0)
            tp = ttf.paragraphs[0]; tp.alignment = PP_ALIGN.CENTER
            tr = tp.add_run(); tr.text = lab_text
            set_jp(tr, size=14, bold=True, color=WHITE)
            # フルラベル (小さく上に)
            _, ftf = add_box(slide, text_x, ry_ - 0.16, txt_inner_w, 0.2)
            fr = ftf.paragraphs[0].add_run(); fr.text = FIELD_FULL[fi]
            set_jp(fr, size=9, color=lab_col)
            # 本文
            _, btf = add_box(slide, text_x, ry_ + 0.05, txt_inner_w, pitch - 0.18)
            val = d[key]
            items = val if isinstance(val, list) else [val]
            for j, t in enumerate(items):
                bp = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
                bp.space_after = Pt(2)
                br = bp.add_run(); br.text = t
                set_jp(br, size=13 if not has_img else 12.5, color=DARK)
            # 区切り線
            if fi < 3:
                add_rect(slide, text_x, ry_ + pitch - 0.12, txt_inner_w, 0.011,
                         RGBColor(0xDD, 0xE0, 0xE4))

        # 画像
        if has_img:
            path = os.path.join(HERE, d["img"])
            sz = png_size(path)
            box_x, box_y, box_w, box_h = 7.75, 1.75, 5.15, 4.55
            if sz:
                iw, ih = sz
                ar = iw / ih
                w = box_w; h = w / ar
                if h > box_h:
                    h = box_h; w = h * ar
                px = box_x + (box_w - w) / 2
                py = box_y + (box_h - h) / 2 - 0.1
                # 枠
                add_rect(slide, px - 0.04, py - 0.04, w + 0.08, h + 0.08,
                         WHITE, line=RGBColor(0xBB, 0xC2, 0xC9), line_w=1.0)
                slide.shapes.add_picture(path, Inches(px), Inches(py),
                                         width=Inches(w), height=Inches(h))
                # キャプション
                _, capf = add_box(slide, box_x - 0.1, py + h + 0.06, box_w + 0.2, 0.7)
                capf.word_wrap = True
                cr = capf.paragraphs[0].add_run()
                cr.text = d.get("img_cap", "")
                capf.paragraphs[0].alignment = PP_ALIGN.CENTER
                set_jp(cr, size=10, italic=True, color=GRAY)
        footer(slide, sec_title, num, TOTAL)

    elif kind == "gap":
        add_rect(slide, 0, 0, 13.333, 0.9, RGBColor(0x10, 0x2A, 0x43))
        _, tf = add_box(slide, 0.45, 0.16, 12.0, 0.6)
        r = tf.paragraphs[0].add_run()
        r.text = "まとめ (1/2)  —  検証された技術空白 = 本ツールの新規性"
        set_jp(r, size=22, bold=True, color=WHITE)
        _, tfb = add_box(slide, 0.6, 1.2, 12.2, 5.7)
        items = [
            ("統合性の空白 (3-0 確証)", "単一の粘性ソルバで「亜音速収縮 + スロート + 超音速コンタ」を一体最適化し、かつ多目的パレートで設計した前例なし。CAN-DO は単目的・NS/PNS 分割、Matsunaga/Ogawa–Boyce/Zhang はコンタのみ・rc 不変。"),
            ("スロート曲率 rc の設計変数化 (3-0 確証)", "rc を自由設計変数として「熱流束 vs 一様性 vs コンパクト性 vs Cd」の多目的パレートで最適化した査読文献は存在しない。最近接の Bianchi 2015 もパラメトリック感度に留まる。"),
            ("凝縮余裕の制約化", "凝縮指標を最適化目的にした実証は湿り蒸気タービンのみ。極超音速風洞の試験ガスでの前例なし。随伴向けの微分可能制約化も未確立。"),
            ("SERN の 3D 統合最適化", "角部 R・サイドフェンス・有限スパンの 3D 効果と、3D・多目的パレート・多成分の統合最適化は文献未確認。"),
            ("結論", "「forge 単一ソルバ × 全域 (rc 含む) × 多目的パレート × 多成分 TP」は検証済み文献における真の空白。本ツールは確立要素技術の組合せで現状技術の最前線 (ないしその先) に位置する。"),
        ]
        for i, (h, t) in enumerate(items):
            p = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
            p.space_before = Pt(10 if i else 0)
            r = p.add_run(); r.text = "▌ " + h
            set_jp(r, size=15, bold=True, color=RGBColor(0xB7, 0x1C, 0x1C))
            p2 = tfb.add_paragraph()
            r2 = p2.add_run(); r2.text = "    " + t
            set_jp(r2, size=13, color=DARK)
        footer(slide, "まとめ", num, TOTAL)

    elif kind == "proposal":
        add_rect(slide, 0, 0, 13.333, 0.9, RGBColor(0x10, 0x2A, 0x43))
        _, tf = add_box(slide, 0.45, 0.16, 12.0, 0.6)
        r = tf.paragraphs[0].add_run()
        r.text = "まとめ (2/2)  —  提案: 3 層アーキテクチャ × 2 フェーズ × 多忠実度"
        set_jp(r, size=22, bold=True, color=WHITE)
        _, tfb = add_box(slide, 0.6, 1.2, 12.2, 5.7)
        items = [
            ("3 層の役割分担 (完全新作にしない)",
             "① 既存 MOC は温存・近代化し「初期形状生成器 + ループ内形状生成器 + 低忠実度層」へ昇格 (目標中心線マッハ分布 Bézier をフロー空間設計変数に)。② forge を粘性 NS 評価器に (軸対称・多成分 TP・k-ω SST・GPU は成熟済)。③ 最適化層のみ新規 (OpenMDAO/DAKOTA 骨格 + pymoo + SMT)。"),
            ("フェーズ A (今すぐ着手・非侵入)",
             "LHS DOE → forge 評価 → Kriging サロゲート → NSGA-II + EGO 的 infill。forge はコード改造ゼロの黒箱接続で、3 機種すべてに即展開可能。スロート曲率 rc・収縮形状も設計変数に開放。"),
            ("フェーズ B (大投資・高効果)",
             "離散随伴の導入。まず SU2 併用 (B-i) で効果を先取りし、必要に応じて forge へ AD 内製 (B-ii, CoDiPack 方式)。随伴勾配は GEK/weighted-GEK でサロゲート精度を底上げし、多目的パレート探索に合流。"),
            ("多忠実度の予算配分",
             "低: MOC/Euler (~秒) = 広域スクリーニング / 中: 軸対称 RANS forge (~分) = パレート本体 / 高: 3D RANS forge (~時間) = infill・最終確認。co-kriging で融合。"),
        ]
        for i, (h, t) in enumerate(items):
            p = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
            p.space_before = Pt(10 if i else 0)
            r = p.add_run(); r.text = "▌ " + h
            set_jp(r, size=15, bold=True, color=RGBColor(0x1B, 0x5E, 0x20))
            p2 = tfb.add_paragraph()
            r2 = p2.add_run(); r2.text = "    " + t
            set_jp(r2, size=13, color=DARK)
        footer(slide, "まとめ", num, TOTAL)

out = os.path.join(HERE, "nozzle_design_survey.pptx")
prs.save(out)
n_img = sum(1 for p in PAPERS if "img" in p and os.path.exists(os.path.join(HERE, p["img"])))
print(f"saved: {out}")
print(f"  slides={TOTAL}  papers={len(PAPERS)}  with-figure={n_img}")
